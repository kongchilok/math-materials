# -*- coding: utf-8 -*-
"""
資源班 SOIL 融合簡報共用建構器（python-pptx）。
配色＝Sage Calm 低刺激版；硬規則見 inclusive-soil-teaching-deck skill 的
references/inclusive-deck-checklist.md D 節：內文≥20pt／標題≥32pt、
每頁≤5行×18字、行距1.3-1.5、淺色底非純白、左對齊、難度只用星星、
無襯線黑體、不用底線斜體、強調用粗體+色塊。
各單元的 build_簡報.py 只放內容資料，版面一律呼叫本模組。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Sage Calm 配色 ----
SAGE = RGBColor(0x84, 0xB5, 0x9F)
EUCALYPTUS = RGBColor(0x69, 0xA2, 0x97)
SLATE = RGBColor(0x50, 0x80, 0x8E)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_BG = RGBColor(0xEA, 0xF2, 0xEE)
CARD_BG = RGBColor(0xD8, 0xE8, 0xE2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WRONG_BG = RGBColor(0xE9, 0xDF, 0xC7)
RIGHT_BG = RGBColor(0xD8, 0xE8, 0xE2)

FONT = "Microsoft JhengHei"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def stars(n):
    return "★" * n + "☆" * (3 - n)


def _set_east_asian_font(run, font_name=FONT):
    rPr = run.font._rPr
    latin = rPr.find(qn('a:latin'))
    if latin is None:
        latin = rPr.makeelement(qn('a:latin'), {})
        rPr.append(latin)
    latin.set('typeface', font_name)
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font_name)


def add_text(slide, left, top, width, height, text, size=20, bold=False,
             color=DARK_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.3, font_name=FONT, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = color
        _set_east_asian_font(run, font_name)
    return box


def add_bullets(slide, left, top, width, height, items, size=20, bold=False,
                 color=DARK_TEXT, line_spacing=1.3, space_after=10,
                 font_name=FONT, marker=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = (marker(i) if marker else "") + item
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = color
        _set_east_asian_font(run, font_name)
    return box


def add_circle_num(slide, cx, cy, d, num, bg=SLATE, fg=WHITE, size=20):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.name = FONT
    run.font.color.rgb = fg
    return shape


def add_card(slide, left, top, width, height, bg=CARD_BG, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_bar(slide, left, top, width, height, bg=EUCALYPTUS, line_color=DARK_TEXT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def add_footer(slide, subject_unit):
    add_text(slide, MARGIN, SLIDE_H - Inches(0.42), SLIDE_W - 2 * MARGIN, Inches(0.32),
              subject_unit, size=11, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.LEFT,
              line_spacing=1.0)


def add_title(slide, text, color=DARK_TEXT, top=MARGIN):
    add_text(slide, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(0.9), text,
              size=32, bold=True, color=color, line_spacing=1.1)


# ========== 頁面角色 ==========

def add_cover(prs, subject, unit_title, mode_label, footer_text):
    s = _blank_slide(prs)
    set_bg(s, SLATE)
    add_text(s, MARGIN, Inches(2.5), SLIDE_W - 2 * MARGIN, Inches(0.6), subject,
              size=22, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s, MARGIN, Inches(3.0), SLIDE_W - 2 * MARGIN, Inches(1.6), unit_title,
              size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT, line_spacing=1.15)
    add_card(s, MARGIN, Inches(4.7), Inches(7.5), Inches(0.9), bg=EUCALYPTUS)
    add_text(s, MARGIN + Inches(0.25), Inches(4.7), Inches(7.0), Inches(0.9), mode_label,
              size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    add_footer_dark(s, footer_text)
    return s


def add_footer_dark(slide, text):
    add_text(slide, MARGIN, SLIDE_H - Inches(0.5), SLIDE_W - 2 * MARGIN, Inches(0.35),
              text, size=11, color=RGBColor(0xDD, 0xE8, 0xE4), line_spacing=1.0)


def add_flow_preview(prs, subject_unit, title, items):
    s = _blank_slide(prs)
    set_bg(s, LIGHT_BG)
    add_title(s, title)
    y = Inches(1.9)
    row_h = Inches(1.0)
    for i, item in enumerate(items):
        add_circle_num(s, MARGIN, y, Inches(0.55), i + 1, bg=SLATE, size=22)
        add_card(s, MARGIN + Inches(0.75), y - Inches(0.05), SLIDE_W - 2 * MARGIN - Inches(0.75), Inches(0.7))
        add_text(s, MARGIN + Inches(1.0), y + Inches(0.08), SLIDE_W - 2 * MARGIN - Inches(1.3), Inches(0.55),
                  item, size=20, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
        y += row_h
    add_footer(s, subject_unit)
    return s


def add_question_intro(prs, subject_unit, title, question, options, hint):
    s = _blank_slide(prs)
    set_bg(s, LIGHT_BG)
    add_title(s, title)
    add_card(s, MARGIN, Inches(1.9), SLIDE_W - 2 * MARGIN, Inches(1.6), bg=CARD_BG)
    add_text(s, MARGIN + Inches(0.3), Inches(2.1), SLIDE_W - 2 * MARGIN - Inches(0.6), Inches(1.2),
              question, size=24, bold=True, line_spacing=1.3, anchor=MSO_ANCHOR.MIDDLE)
    y = Inches(3.8)
    for i, opt in enumerate(options):
        add_circle_num(s, MARGIN, y, Inches(0.5), chr(65 + i), bg=EUCALYPTUS, size=20)
        add_text(s, MARGIN + Inches(0.7), y + Inches(0.05), SLIDE_W - 2 * MARGIN - Inches(0.7), Inches(0.5),
                  opt, size=20, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.75)
    add_text(s, MARGIN, Inches(6.4), SLIDE_W - 2 * MARGIN, Inches(0.6), hint,
              size=20, color=RGBColor(0x44, 0x44, 0x44), line_spacing=1.2)
    add_footer(s, subject_unit)
    return s


def add_cra(prs, subject_unit, title, concrete, representational, abstract):
    """concrete/representational/abstract: dict(label=, desc=)"""
    s = _blank_slide(prs)
    set_bg(s, LIGHT_BG)
    add_title(s, title)
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) / 3
    tops = Inches(1.9)
    cols = [concrete, representational, abstract]
    labels_bg = [SAGE, EUCALYPTUS, SLATE]
    for i, col in enumerate(cols):
        x = MARGIN + i * (col_w + Inches(0.3))
        add_card(s, x, tops, col_w, Inches(0.6), bg=labels_bg[i])
        add_text(s, x, tops, col_w, Inches(0.6), col["label"], size=20, bold=True,
                  color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_card(s, x, tops + Inches(0.75), col_w, Inches(3.9), bg=CARD_BG)
        add_text(s, x + Inches(0.2), tops + Inches(0.95), col_w - Inches(0.4), Inches(3.5),
                  col["desc"], size=20, line_spacing=1.3, anchor=MSO_ANCHOR.TOP)
    add_footer(s, subject_unit)
    return s


def add_step_card(prs, subject_unit, title, steps, note=None):
    """steps: list[str]. Row height adapts to len(steps) so long lists never
    push the note/footer past the slide bottom (QB: unit19 6-item overflow)."""
    s = _blank_slide(prs)
    set_bg(s, LIGHT_BG)
    add_title(s, title)
    y0 = Inches(1.9)
    footer_top = SLIDE_H - Inches(0.55)
    note_reserve = Inches(0.8) if note else Inches(0)
    available = footer_top - y0 - note_reserve
    ideal_row_h = Emu(int(available / max(len(steps), 1)))
    row_h = min(Inches(0.85), max(Inches(0.55), ideal_row_h))
    y = y0
    text_size = 20 if row_h >= Inches(0.65) else 16
    for i, step in enumerate(steps):
        circle_d = min(Inches(0.5), row_h - Inches(0.05))
        add_circle_num(s, MARGIN, y, circle_d, i + 1, bg=SLATE, size=min(20, text_size))
        add_text(s, MARGIN + Inches(0.75), y + Inches(0.02), SLIDE_W - 2 * MARGIN - Inches(0.75),
                  row_h - Inches(0.05), step, size=text_size, bold=True, anchor=MSO_ANCHOR.MIDDLE,
                  line_spacing=1.15)
        y += row_h
    if note:
        add_card(s, MARGIN, y + Inches(0.15), SLIDE_W - 2 * MARGIN, Inches(0.65), bg=CARD_BG)
        add_text(s, MARGIN + Inches(0.25), y + Inches(0.27), SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.45),
                  note, size=min(20, text_size), anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, subject_unit)
    return s


def add_transition(prs, text, footer_text):
    s = _blank_slide(prs)
    set_bg(s, SLATE)
    add_text(s, Inches(1.2), Inches(3.0), SLIDE_W - Inches(2.4), Inches(1.6), text,
              size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
              line_spacing=1.3)
    add_footer_dark(s, footer_text)
    return s


def add_misconception(prs, subject_unit, title, items):
    """items: list of dict(wrong=, right=)"""
    s = _blank_slide(prs)
    set_bg(s, LIGHT_BG)
    add_title(s, title)
    y = Inches(1.9)
    row_h = Inches(1.5)
    for item in items:
        add_card(s, MARGIN, y, (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2, Inches(1.25), bg=WRONG_BG)
        add_text(s, MARGIN + Inches(0.2), y + Inches(0.1), (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2 - Inches(0.4),
                  Inches(1.05), "✗　" + item["wrong"], size=20, bold=True, line_spacing=1.2,
                  anchor=MSO_ANCHOR.MIDDLE)
        x2 = MARGIN + (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2 + Inches(0.3)
        add_card(s, x2, y, (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2, Inches(1.25), bg=RIGHT_BG)
        add_text(s, x2 + Inches(0.2), y + Inches(0.1), (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2 - Inches(0.4),
                  Inches(1.05), "✓　" + item["right"], size=20, bold=True, line_spacing=1.2,
                  anchor=MSO_ANCHOR.MIDDLE)
        y += row_h
    add_footer(s, subject_unit)
    return s


def add_scaffold_2x2(prs, subject_unit, title, problem, what, given, strategy, check):
    s = _blank_slide(prs)
    set_bg(s, LIGHT_BG)
    add_title(s, title)
    add_card(s, MARGIN, Inches(1.6), SLIDE_W - 2 * MARGIN, Inches(0.65), bg=CARD_BG)
    add_text(s, MARGIN + Inches(0.25), Inches(1.6), SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.65),
              problem, size=18, bold=True, line_spacing=1.1, anchor=MSO_ANCHOR.MIDDLE)
    cells = [("問什麼", what, SAGE), ("已知", given, EUCALYPTUS),
             ("策略算式", strategy, SLATE), ("解答驗算", check, SAGE)]
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    row_h = Inches(2.15)
    top0 = Inches(2.4)
    for i, (label, text, bg) in enumerate(cells):
        col = i % 2
        row = i // 2
        x = MARGIN + col * (col_w + Inches(0.3))
        y = top0 + row * (row_h + Inches(0.15))
        add_card(s, x, y, col_w, Inches(0.45), bg=bg)
        add_text(s, x, y, col_w, Inches(0.45), label, size=18, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_card(s, x, y + Inches(0.5), col_w, row_h - Inches(0.5), bg=CARD_BG)
        add_text(s, x + Inches(0.2), y + Inches(0.6), col_w - Inches(0.4), row_h - Inches(0.7),
                  text, size=20, line_spacing=1.15)
    add_footer(s, subject_unit)
    return s


def add_practice_explain(prs, subject_unit, title, problem, steps, note=None):
    s = _blank_slide(prs)
    set_bg(s, LIGHT_BG)
    add_title(s, title)
    add_card(s, MARGIN, Inches(1.75), SLIDE_W - 2 * MARGIN, Inches(0.8), bg=CARD_BG)
    add_text(s, MARGIN + Inches(0.25), Inches(1.85), SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.6),
              problem, size=20, bold=True, line_spacing=1.15, anchor=MSO_ANCHOR.MIDDLE)
    y = Inches(2.8)
    for i, step in enumerate(steps):
        add_circle_num(s, MARGIN, y, Inches(0.44), i + 1, bg=SLATE, size=18)
        add_text(s, MARGIN + Inches(0.62), y - Inches(0.02), SLIDE_W - 2 * MARGIN - Inches(0.62), Inches(0.5),
                  step, size=20, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
        y += Inches(0.62)
    if note:
        add_card(s, MARGIN, y + Inches(0.1), SLIDE_W - 2 * MARGIN, Inches(0.6), bg=WRONG_BG)
        add_text(s, MARGIN + Inches(0.25), y + Inches(0.2), SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.4),
                  note, size=18, bold=True)
    add_footer(s, subject_unit)
    return s


def add_tiered_task(prs, subject_unit, title, tasks):
    """tasks: list of dict(star=n, label=, problem=)"""
    s = _blank_slide(prs)
    set_bg(s, LIGHT_BG)
    add_title(s, title)
    n = len(tasks)
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3) * (n - 1)) / n
    colors = [SAGE, EUCALYPTUS, SLATE]
    for i, task in enumerate(tasks):
        x = MARGIN + i * (col_w + Inches(0.3))
        add_card(s, x, Inches(1.9), col_w, Inches(0.7), bg=colors[i % 3])
        add_text(s, x, Inches(1.9), col_w, Inches(0.35), task["label"], size=18, bold=True,
                  color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.22), col_w, Inches(0.35), stars(task["star"]), size=16,
                  color=WHITE, align=PP_ALIGN.CENTER)
        add_card(s, x, Inches(2.75), col_w, Inches(4.0), bg=CARD_BG)
        add_text(s, x + Inches(0.18), Inches(2.9), col_w - Inches(0.36), Inches(3.7),
                  task["problem"], size=20, line_spacing=1.2)
    add_footer(s, subject_unit)
    return s


def add_summary(prs, takeaway, floor_version, footer_text):
    s = _blank_slide(prs)
    set_bg(s, SLATE)
    add_text(s, MARGIN, Inches(1.4), SLIDE_W - 2 * MARGIN, Inches(0.6), "帶走一句話",
              size=20, bold=True, color=RGBColor(0xDD, 0xE8, 0xE4))
    add_text(s, MARGIN, Inches(2.0), SLIDE_W - 2 * MARGIN, Inches(1.5), takeaway,
              size=34, bold=True, color=WHITE, line_spacing=1.3)
    add_card(s, MARGIN, Inches(4.1), SLIDE_W - 2 * MARGIN, Inches(1.7), bg=EUCALYPTUS)
    add_text(s, MARGIN + Inches(0.3), Inches(4.3), SLIDE_W - 2 * MARGIN - Inches(0.6), Inches(0.4),
              "保底版（至少要能說出的一句）", size=18, bold=True, color=WHITE)
    add_text(s, MARGIN + Inches(0.3), Inches(4.75), SLIDE_W - 2 * MARGIN - Inches(0.6), Inches(0.9),
              floor_version, size=22, bold=True, color=WHITE, line_spacing=1.25)
    add_footer_dark(s, footer_text)
    return s


def save(prs, path):
    prs.save(path)
